"""
PowerPoint presentation generator from video frames.
"""

from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import io
from typing import List, Dict


class PPTXGenerator:
    """Generates PowerPoint presentations from video frames."""
    
    def __init__(self, slide_width: float = 10.0, slide_height: float = 7.5):
        """
        Initialize PPTX generator.
        
        Args:
            slide_width: Slide width in inches (default 10.0 for 16:9 aspect ratio)
            slide_height: Slide height in inches (default 7.5 for 16:9 aspect ratio)
        """
        self.slide_width = slide_width
        self.slide_height = slide_height
        self.max_image_width = Inches(slide_width - 0.5)  # Leave margins
        self.max_image_height = Inches(slide_height - 0.5)
    
    def create_presentation(self, frames: List[Dict], output_path: str = None) -> bytes:
        """
        Create PowerPoint presentation from selected frames.
        
        Args:
            frames: List of frame dictionaries with 'image' PIL Image objects
            output_path: Optional path to save file. If None, returns bytes.
            
        Returns:
            Bytes of PPTX file if output_path is None, otherwise None
        """
        # Create presentation
        prs = Presentation()
        
        # Set slide size (16:9 aspect ratio)
        prs.slide_width = Inches(self.slide_width)
        prs.slide_height = Inches(self.slide_height)
        
        print(f"Creating presentation with {len(frames)} slides...")
        
        for idx, frame_data in enumerate(frames):
            # Create blank slide layout
            blank_slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Get image and ensure it's a PIL Image
            image = frame_data['image']
            
            # Validate and convert to PIL Image if needed
            if not isinstance(image, Image.Image):
                if hasattr(image, 'read'):
                    # It's a file-like object (BytesIO, file handle, etc.)
                    if hasattr(image, 'seek'):
                        image.seek(0)
                    image_bytes = image.read()
                    if not image_bytes:
                        raise ValueError(f"Frame {idx} has empty image data")
                    image = Image.open(io.BytesIO(image_bytes))
                else:
                    raise ValueError(f"Frame {idx} has invalid image type: {type(image)}")
            
            # Calculate dimensions maintaining aspect ratio
            img_width, img_height = image.size
            aspect_ratio = img_width / img_height
            
            # Fit image within slide bounds while maintaining aspect ratio
            if aspect_ratio > (self.max_image_width / self.max_image_height):
                # Image is wider - fit to width
                width = self.max_image_width
                height = width / aspect_ratio
            else:
                # Image is taller - fit to height
                height = self.max_image_height
                width = height * aspect_ratio
            
            # Center image on slide
            left = (Inches(self.slide_width) - width) / 2
            top = (Inches(self.slide_height) - height) / 2
            
            # Convert PIL Image to bytes for python-pptx
            img_bytes = io.BytesIO()
            # Resize image if it's too large (optimize for PPTX)
            if img_width > 1920 or img_height > 1080:
                # Resize to max 1920x1080 while maintaining aspect ratio
                max_dimension = 1920
                if img_width > img_height:
                    new_width = max_dimension
                    new_height = int(img_height * (max_dimension / img_width))
                else:
                    new_height = max_dimension
                    new_width = int(img_width * (max_dimension / img_height))
                image = image.resize((new_width, new_height), Image.Resampling.LANCZOS)
            
            image.save(img_bytes, format='PNG')
            img_bytes.seek(0)
            
            # Add image to slide
            # python-pptx's add_picture accepts BytesIO, but ensure it's in the right state
            # If it fails, we'll get a clear error message
            try:
                slide.shapes.add_picture(img_bytes, left, top, width, height)
            except (TypeError, AttributeError) as e:
                # If BytesIO doesn't work, create a fresh one with proper state
                img_bytes.seek(0)
                image_data = img_bytes.getvalue()
                fresh_img_bytes = io.BytesIO(image_data)
                slide.shapes.add_picture(fresh_img_bytes, left, top, width, height)
            
            if (idx + 1) % 10 == 0:
                print(f"Added {idx + 1}/{len(frames)} slides...")
        
        print(f"Presentation created successfully with {len(frames)} slides")
        
        # Save to file or return bytes
        if output_path:
            prs.save(output_path)
            return None
        else:
            # Save to bytes
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            return output.getvalue()
    
    def create_from_selected_indices(self, all_frames: List[Dict], 
                                     selected_indices: List[int],
                                     output_path: str = None) -> bytes:
        """
        Create presentation from selected frame indices.
        
        Args:
            all_frames: All available frames
            selected_indices: List of frame indices to include
            output_path: Optional path to save file
            
        Returns:
            Bytes of PPTX file if output_path is None, otherwise None
        """
        selected_frames = [all_frames[i] for i in selected_indices if 0 <= i < len(all_frames)]
        return self.create_presentation(selected_frames, output_path)
